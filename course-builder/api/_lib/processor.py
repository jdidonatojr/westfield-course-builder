"""
AILA Course Builder — PPTX Processing Engine

Takes a .pptx file and the editorial info (Course Type, Audience,
Section Map, Teaching Notes), then produces:

  1. The instructor notes .txt file
  2. The ElevenLabs catalog snippet
  3. The Yola page snippet
  4. The assembly checklist

PDF generation happens separately via CloudConvert.
"""

from pptx import Presentation
import re


# ============================================================
# HELPER: turn a course title into a safe filename
# ============================================================
def make_filename(course_title):
    """
    Convert 'Ace Your Next Job Interview' to 'Ace-Your-Next-Job-Interview'.

    Removes punctuation that breaks filenames and replaces spaces
    with hyphens. Keeps it readable for humans.
    """
    # Remove characters that aren't letters, numbers, spaces, or hyphens
    cleaned = re.sub(r'[^\w\s-]', '', course_title)
    # Replace spaces with hyphens
    cleaned = re.sub(r'\s+', '-', cleaned.strip())
    return cleaned


# ============================================================
# HELPER: turn a course title into a short slug for course IDs
# ============================================================
def make_slug(course_title):
    """
    Convert 'Ace Your Next Job Interview' to 'ace_your_next_job_interview'.

    Used for the Yola course ID and ElevenLabs internal references.
    """
    cleaned = re.sub(r'[^\w\s-]', '', course_title)
    cleaned = re.sub(r'\s+', '_', cleaned.strip()).lower()
    return cleaned


# ============================================================
# HELPER: get visible text from a slide (excludes speaker notes)
# ============================================================
def get_visible_text(slide):
    """Pull all visible text from a slide's shapes."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                parts.append(txt)
    return '\n'.join(parts)


# ============================================================
# HELPER: get speaker notes from a slide
# ============================================================
def get_speaker_notes(slide):
    """Pull the speaker notes from a slide, if any exist."""
    if slide.has_notes_slide:
        return slide.notes_slide.notes_text_frame.text.strip()
    return ""


# ============================================================
# HELPER: check if a slide is hidden
# ============================================================
def is_hidden(slide):
    """
    Detect if a slide is marked as hidden in PowerPoint.
    Hidden slides are skipped during the slideshow.
    """
    # python-pptx exposes this via the show attribute on the XML element
    show = slide._element.get('show')
    # show='0' means hidden; absence of attribute means visible
    return show == '0'


# ============================================================
# MAIN FUNCTION: process the deck and return all the pieces
# ============================================================
def process_deck(pptx_path, course_info):
    """
    The main entry point. Reads a .pptx and returns a dictionary
    with all the files and content needed to publish the course.

    Args:
        pptx_path: path to the .pptx file
        course_info: dictionary with these required keys:
            - course_number   (e.g., 5)
            - course_title    (e.g., "Ace Your Next Job Interview")
            - course_type     (e.g., "Skills-based, hands-on")
            - audience        (e.g., "Job seekers preparing for interviews")
            - section_map     (list of strings, one per section)
            - teaching_notes  (list of strings, one per bullet)

    Returns:
        Dictionary with:
            - filename_base   (e.g., "Ace-Your-Next-Job-Interview")
            - pdf_filename    (e.g., "Ace-Your-Next-Job-Interview.pdf")
            - txt_filename    (e.g., "Ace-Your-Next-Job-Interview.txt")
            - txt_content     (the full instructor notes text)
            - elevenlabs_snippet  (the catalog block)
            - yola_snippet    (the JavaScript object entry)
            - checklist       (the assembly instructions in Markdown)
            - total_slides    (count of visible slides)
            - skipped_slides  (count of hidden slides we left out)
    """

    # ----------- STEP 1: read the deck -----------
    presentation = Presentation(pptx_path)

    # Separate visible from hidden slides
    visible_slides = []
    hidden_count = 0
    for slide in presentation.slides:
        if is_hidden(slide):
            hidden_count += 1
        else:
            visible_slides.append(slide)

    total_visible = len(visible_slides)

    # ----------- STEP 2: build filenames -----------
    filename_base = make_filename(course_info['course_title'])
    pdf_filename = f"{filename_base}.pdf"
    txt_filename = f"{filename_base}.txt"

    # ----------- STEP 3: build the .txt content -----------
    txt_content = build_txt_file(
        course_info['course_title'],
        visible_slides
    )

    # ----------- STEP 4: build the ElevenLabs snippet -----------
    elevenlabs_snippet = build_elevenlabs_snippet(
        course_info,
        total_visible,
        pdf_filename,
        txt_filename
    )

    # ----------- STEP 5: build the Yola snippet -----------
    yola_snippet = build_yola_snippet(
        course_info,
        total_visible,
        pdf_filename
    )

    # ----------- STEP 6: build the assembly checklist -----------
    checklist = build_checklist(
        course_info,
        pdf_filename,
        txt_filename
    )

    # ----------- STEP 7: return everything -----------
    return {
        'filename_base': filename_base,
        'pdf_filename': pdf_filename,
        'txt_filename': txt_filename,
        'txt_content': txt_content,
        'elevenlabs_snippet': elevenlabs_snippet,
        'yola_snippet': yola_snippet,
        'checklist': checklist,
        'total_slides': total_visible,
        'skipped_slides': hidden_count
    }


# ============================================================
# BUILDER: the instructor notes .txt file
# ============================================================
def build_txt_file(course_title, visible_slides):
    """
    Builds the full instructor notes file in the exact format
    Joe established with the Reclaim Your Time sample.
    """
    lines = []
    lines.append(course_title.upper())
    lines.append("Full course content for AI instructor")
    lines.append("=" * 70)
    lines.append("")
    lines.append("")

    for i, slide in enumerate(visible_slides, 1):
        lines.append(f"=== SLIDE {i} ===")
        lines.append("")
        lines.append("Visible content on slide:")
        visible = get_visible_text(slide)
        lines.append(visible if visible else "(no visible text)")
        lines.append("")
        lines.append("Speaker notes (how to teach this slide):")
        notes = get_speaker_notes(slide)
        lines.append(notes if notes else "(no speaker notes)")
        lines.append("")
        lines.append("-" * 70)
        lines.append("")

    return '\n'.join(lines)


# ============================================================
# BUILDER: same as above, but from pre-extracted slide data (browser-parsed)
# ============================================================
def build_txt_file_from_data(course_title, slides_data):
    """
    Build the TXT file from slide data that was extracted in the browser.
    Each item in slides_data is a dict with 'visible' and 'notes' keys.
    """
    lines = []
    lines.append(course_title.upper())
    lines.append("Full course content for AI instructor")
    lines.append("=" * 70)
    lines.append("")
    lines.append("")

    for i, slide in enumerate(slides_data, 1):
        lines.append(f"=== SLIDE {i} ===")
        lines.append("")
        lines.append("Visible content on slide:")
        visible = slide.get('visible', '').strip()
        lines.append(visible if visible else "(no visible text)")
        lines.append("")
        lines.append("Speaker notes (how to teach this slide):")
        notes = slide.get('notes', '').strip()
        lines.append(notes if notes else "(no speaker notes)")
        lines.append("")
        lines.append("-" * 70)
        lines.append("")

    return '\n'.join(lines)


# ============================================================
# BUILDER: the ElevenLabs catalog snippet
# ============================================================
def build_elevenlabs_snippet(course_info, total_slides, pdf_filename, txt_filename):
    """
    Builds the catalog block that goes into the ElevenLabs
    system prompt. Matches the format of Courses #1-#4.
    """
    course_num = course_info['course_number']
    title = course_info['course_title']

    lines = []
    lines.append("# -------------------------------------------------------")
    lines.append(f"# COURSE #{course_num} — {title}")
    lines.append("# -------------------------------------------------------")
    lines.append(f"Course Title: {title}")
    lines.append(f"Total Slides: {total_slides}")
    lines.append(f"Slide Content File (PDF): {pdf_filename}")
    lines.append(f"Teaching Notes File (TXT): {txt_filename}")
    lines.append(f"Course Type: {course_info['course_type']}")
    lines.append(f"Audience: {course_info['audience']}")
    lines.append("SECTION MAP:")
    for section in course_info['section_map']:
        lines.append(f"- {section}")
    lines.append("TEACHING NOTES:")
    for note in course_info['teaching_notes']:
        lines.append(f"- {note}")

    return '\n'.join(lines)


# ============================================================
# BUILDER: the Yola JavaScript snippet
# ============================================================
def build_yola_snippet(course_info, total_slides, pdf_filename):
    """
    Builds the JavaScript object entry that goes into the
    COURSE_LIBRARY in the Yola page code.
    """
    course_num = course_info['course_number']
    title = course_info['course_title']
    slug = make_slug(title)

    # Build section entries for the JS object
    # We need to extract slide numbers from section labels like
    # "Introduction (slides 1-4)" → startSlide: 1
    section_entries = []
    for section in course_info['section_map']:
        start = extract_start_slide(section)
        # Strip the slide range from the label for cleaner display
        clean_label = section
        section_entries.append(
            f"        {{ label: '{escape_js(clean_label)}', startSlide: {start} }}"
        )

    github_base = "https://raw.githubusercontent.com/jdidonatojr/aila-course-files/main"

    lines = []
    lines.append(f'    "{course_num}": {{')
    lines.append(f"      id: '{slug}',")
    lines.append(f"      title: '{escape_js(title)}',")
    lines.append(f"      totalSlides: {total_slides},")
    lines.append(f"      pdfUrl: '{github_base}/{pdf_filename}',")
    lines.append(f"      firstMessage: 'Welcome to {escape_js(title)}. Are you ready to begin?',")
    lines.append(f"      sections: [")
    lines.append(',\n'.join(section_entries))
    lines.append(f"      ]")
    lines.append(f"    }}")

    return '\n'.join(lines)


# ============================================================
# HELPER: pull the starting slide number from a section label
# ============================================================
def extract_start_slide(section_label):
    """
    From 'Introduction (slides 1-4)' extract 1.
    From 'Module 3 — Countering (16-22)' extract 16.
    Falls back to 1 if no number found.
    """
    # Look for the first number in parentheses or after "slides"
    match = re.search(r'\((?:slides?\s+)?(\d+)', section_label, re.IGNORECASE)
    if match:
        return int(match.group(1))
    # Fallback: first number anywhere in the string
    match = re.search(r'\d+', section_label)
    if match:
        return int(match.group(0))
    return 1


# ============================================================
# HELPER: escape single quotes for JavaScript strings
# ============================================================
def escape_js(text):
    """Make text safe to embed inside JavaScript single-quoted strings."""
    return text.replace("\\", "\\\\").replace("'", "\\'")


# ============================================================
# BUILDER: the assembly checklist
# ============================================================
def build_checklist(course_info, pdf_filename, txt_filename):
    """
    Builds a step-by-step Markdown checklist telling the team
    exactly where each file goes and what to paste where.
    """
    course_num = course_info['course_number']
    title = course_info['course_title']

    checklist = f"""# Assembly Checklist — Course #{course_num}: {title}

This checklist tells you exactly where to put each file and snippet
from this ZIP package. Work through each step in order.

---

## Step 1 — Upload PDF to GitHub

The PDF powers the slide carousel on the Yola page.

- [ ] Go to the `aila-course-files` repo on GitHub
- [ ] Upload `{pdf_filename}` to the root of the repo
- [ ] Commit with message: "Add Course #{course_num} PDF"

---

## Step 2 — Upload PDF and TXT to ElevenLabs

The agent uses these files to teach the course.

- [ ] Open your ElevenLabs agent
- [ ] Go to the Knowledge Base section
- [ ] Upload `{pdf_filename}`
- [ ] Upload `{txt_filename}`

---

## Step 3 — Update the ElevenLabs system prompt

Add the new course to the catalog block.

- [ ] Open the agent's system prompt
- [ ] Find the `# COURSE LIBRARY` section near the bottom
- [ ] Paste the content of `elevenlabs-snippet.txt` after the last existing course
- [ ] Save the agent

---

## Step 4 — Update the Yola page

Add the new course to the course library.

- [ ] Open your Yola page editor
- [ ] Find the `COURSE_LIBRARY` JavaScript object in the page code
- [ ] Paste the content of `yola-snippet.txt` as a new entry
- [ ] Add a comma after the previous course's closing brace if needed
- [ ] Publish the page

---

## Step 5 — Test the course

Before announcing, verify everything works.

- [ ] Visit your Yola catalog page
- [ ] Click the new course
- [ ] Confirm the agent greets you with the right course name
- [ ] Click through 2-3 slides to confirm the carousel works
- [ ] Ask the agent a test question to confirm the knowledge base loaded

---

## Done! 🎉

If anything doesn't work, the most common issues are:
- PDF not uploaded to GitHub root (slide carousel shows blank)
- ElevenLabs knowledge base files not uploaded (agent says "I don't have that course")
- Yola JavaScript has a syntax error (page shows error screen)
"""
    return checklist
